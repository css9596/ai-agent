"""다양한 파일 형식을 텍스트로 변환하는 유틸리티"""

from pathlib import Path
from typing import Optional, Tuple
from dataclasses import dataclass

from PyPDF2 import PdfReader
from openpyxl import load_workbook
from pptx import Presentation


@dataclass
class ExtractionResult:
    """파일 추출 결과"""
    success: bool
    text: str = ""
    error: Optional[str] = None
    file_size: int = 0
    page_count: int = 0  # PDF 페이지 수 또는 Sheet 수


def extract_text_from_file(file_path: str, file_name: str) -> ExtractionResult:
    """파일 형식에 따라 텍스트 추출 (개선된 버전)"""
    file_ext = Path(file_name).suffix.lower()

    try:
        file_size = Path(file_path).stat().st_size

        if file_ext == ".pdf":
            result = extract_from_pdf(file_path)
        elif file_ext in [".xlsx", ".xls"]:
            result = extract_from_excel(file_path)
        elif file_ext in [".pptx", ".ppt"]:
            result = extract_from_powerpoint(file_path)
        elif file_ext in [".txt", ".md"]:
            result = read_text_file(file_path)
        else:
            result = read_text_file(file_path)

        result.file_size = file_size
        return result

    except Exception as e:
        return ExtractionResult(
            success=False,
            error=f"파일 처리 중 오류 발생: {str(e)}"
        )


def read_text_file(file_path: str) -> ExtractionResult:
    """일반 텍스트 파일 읽기"""
    try:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        except UnicodeDecodeError:
            # UTF-8 실패 시 다른 인코딩 시도
            with open(file_path, "r", encoding="euc-kr") as f:
                text = f.read()

        return ExtractionResult(success=True, text=text, page_count=1)

    except Exception as e:
        return ExtractionResult(
            success=False,
            error=f"텍스트 파일 읽기 실패: {str(e)}"
        )


def extract_from_pdf(file_path: str) -> ExtractionResult:
    """PDF에서 텍스트 추출"""
    text = []
    page_count = 0

    try:
        with open(file_path, "rb") as f:
            reader = PdfReader(f)
            page_count = len(reader.pages)

            for page_num, page in enumerate(reader.pages, 1):
                text.append(f"--- Page {page_num} ---")
                try:
                    extracted = page.extract_text()
                    if extracted:
                        text.append(extracted)
                except Exception as e:
                    text.append(f"[페이지 {page_num} 추출 실패: {str(e)}]")

        return ExtractionResult(
            success=True,
            text="\n".join(text),
            page_count=page_count
        )

    except Exception as e:
        return ExtractionResult(
            success=False,
            error=f"PDF 읽기 실패: {str(e)}"
        )


def extract_from_excel(file_path: str) -> ExtractionResult:
    """Excel에서 텍스트 추출"""
    text = []
    sheet_count = 0

    try:
        wb = load_workbook(file_path)
        sheet_count = len(wb.sheetnames)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text.append(f"\n=== Sheet: {sheet_name} ===")

            # 헤더 추출
            headers = []
            for cell in ws[1]:
                headers.append(str(cell.value) if cell.value is not None else "")

            if headers and any(headers):  # 헤더가 있으면
                text.append(" | ".join(headers))
                text.append("-" * 50)

            # 데이터 추출
            for row in ws.iter_rows(min_row=2, values_only=True):
                row_str = " | ".join(str(v) if v is not None else "" for v in row)
                if row_str.strip():  # 빈 줄 제외
                    text.append(row_str)

        return ExtractionResult(
            success=True,
            text="\n".join(text),
            page_count=sheet_count
        )

    except Exception as e:
        return ExtractionResult(
            success=False,
            error=f"Excel 읽기 실패: {str(e)}"
        )


def extract_from_powerpoint(file_path: str) -> ExtractionResult:
    """PowerPoint에서 텍스트 추출"""
    text = []
    slide_count = 0

    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, 1):
            text.append(f"\n=== Slide {slide_num} ===")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)

                # 테이블 처리
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip():
                            text.append(row_text)

        return ExtractionResult(
            success=True,
            text="\n".join(text),
            page_count=slide_count
        )

    except Exception as e:
        return ExtractionResult(
            success=False,
            error=f"PowerPoint 읽기 실패: {str(e)}"
        )


def is_supported_file(file_name: str) -> bool:
    """지원하는 파일 형식인지 확인"""
    supported_extensions = {
        ".pdf",
        ".txt",
        ".md",
        ".xlsx",
        ".xls",
        ".pptx",
        ".ppt",
    }
    return Path(file_name).suffix.lower() in supported_extensions


def get_supported_formats() -> list:
    """지원 형식 목록"""
    return [".pdf", ".txt", ".md", ".xlsx", ".xls", ".pptx", ".ppt"]
